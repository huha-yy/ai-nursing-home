#!/usr/bin/env python3
"""
Markdown → PPTX 生成器 —— 将 markdown 文章/大纲转换为 PPT。

用法:
  python generate_ppt.py --input article.md --output output.pptx

支持从 markdown 标题层级自动分页：
  # 标题     → 封面页
  ## 章节   → 章节标题页 + 内容页
  ### 子节  → 内容点

可选参数:
  --theme <科技蓝|商务|简约>  配色方案（默认 科技蓝）
  --brand-name <品牌名>       品牌名（显示在封面副标题）
"""

import os
import sys
import re
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 配色方案 ────────────────────────────────────────────────────

THEMES = {
    "科技蓝": {
        "primary": RGBColor(0x1A, 0x56, 0xDB),      # 主色蓝
        "secondary": RGBColor(0x2D, 0x37, 0x48),     # 深色
        "accent": RGBColor(0x00, 0xB8, 0xD4),        # 亮蓝
        "bg": RGBColor(0xF8, 0xFA, 0xFC),            # 浅灰白
        "text": RGBColor(0x33, 0x33, 0x33),           # 正文深灰
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),     # 标题白
    },
    "商务": {
        "primary": RGBColor(0x1F, 0x2A, 0x44),
        "secondary": RGBColor(0x3A, 0x4A, 0x6B),
        "accent": RGBColor(0xC9, 0x96, 0x5A),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x2C, 0x2C, 0x2C),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "简约": {
        "primary": RGBColor(0x2D, 0x2D, 0x2D),
        "secondary": RGBColor(0x5A, 0x5A, 0x5A),
        "accent": RGBColor(0xE8, 0x4D, 0x3D),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x33, 0x33, 0x33),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
    },
}

DEFAULT_THEME = "科技蓝"
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)


def parse_markdown(filepath: str) -> dict:
    """解析 markdown 文件，返回结构化内容。"""
    with open(filepath, encoding="utf-8") as f:
        content = f.read()

    lines = content.split("\n")
    title = ""
    sections = []
    current_section = None
    current_subsections = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# ") and not title:
            title = stripped.lstrip("# ").strip()
        elif stripped.startswith("## "):
            if current_section:
                current_section["subsections"] = current_subsections
                sections.append(current_section)
            current_section = {
                "title": stripped.lstrip("## ").strip(),
                "content": [],
                "subsections": [],
            }
            current_subsections = []
        elif stripped.startswith("### "):
            if current_section:
                if current_section["content"]:
                    current_section["subsections"].append({
                        "title": "",
                        "content": current_section["content"][:],
                    })
                    current_section["content"] = []
                current_subsections.append({
                    "title": stripped.lstrip("### ").strip(),
                    "content": [],
                })
        elif current_section:
            if stripped:
                target = current_subsections[-1]["content"] if current_subsections else current_section["content"]
                # Clean markdown formatting for slide text
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
                clean = re.sub(r"\*(.+?)\*", r"\1", clean)
                target.append(clean)

    if current_section:
        current_section["subsections"] = current_subsections
        sections.append(current_section)

    return {"title": title or "未命名文档", "sections": sections}


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """在 slide 上添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or RGBColor(0x33, 0x33, 0x33)
    p.alignment = alignment
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, items, font_size=16,
                        color=None, spacing=Pt(6)):
    """添加带要点的文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Handle bullet/numbered list detection
        if item.startswith("- ") or item.startswith("* "):
            p.text = item[2:]
            p.level = 0
        elif re.match(r"^\d+[\.\、]", item):
            p.text = item
        else:
            p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or RGBColor(0x33, 0x33, 0x33)
        p.space_after = spacing

    return txBox


def _add_cover_slide(prs, title, theme_colors, brand_name=""):
    """创建封面页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 背景色块（上半）
    bg_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(4.5),
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = theme_colors["primary"]
    bg_shape.line.fill.background()

    # 标题
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(2),
                 title, font_size=40, color=theme_colors["title_text"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # 副标题
    subtitle = brand_name or ""
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                     subtitle, font_size=20, color=RGBColor(0xDD, 0xDD, 0xDD),
                     alignment=PP_ALIGN.LEFT)

    return slide


def _add_section_slide(prs, section_title, section_num, total, theme_colors):
    """创建章节过渡页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 左侧色块
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.3), SLIDE_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme_colors["accent"]
    bar.line.fill.background()

    # 章节编号
    _add_textbox(slide, Inches(1), Inches(2), Inches(2), Inches(0.8),
                 f"0{section_num}" if section_num < 10 else str(section_num),
                 font_size=48, color=theme_colors["accent"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # 章节标题
    _add_textbox(slide, Inches(1), Inches(3), Inches(10), Inches(1.5),
                 section_title, font_size=32, color=theme_colors["secondary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    return slide


def _add_content_slide(prs, title, items, theme_colors, is_sub=False):
    """创建内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 顶部标题条
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2),
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = theme_colors["secondary"]
    title_bar.line.fill.background()

    _add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8),
                 title, font_size=24, color=RGBColor(0xFF, 0xFF, 0xFF),
                 bold=True, alignment=PP_ALIGN.LEFT)

    if items:
        _add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(5.5),
                            items, font_size=15 if is_sub else 16,
                            color=theme_colors["text"])

    return slide


def generate_ppt(input_path: str, output_path: str, theme_name: str = DEFAULT_THEME,
                 brand_name: str = "") -> str:
    """从 markdown 生成 PPTX 文件。"""
    theme_name = theme_name if theme_name in THEMES else DEFAULT_THEME
    theme_colors = THEMES[theme_name]

    doc = parse_markdown(input_path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1) 封面
    _add_cover_slide(prs, doc["title"], theme_colors, brand_name)

    # 2) 各章节
    total = len(doc["sections"])
    for i, section in enumerate(doc["sections"], 1):
        # 章节过渡页
        _add_section_slide(prs, section["title"], i, total, theme_colors)

        # 内容页
        if section["content"]:
            _add_content_slide(prs, section["title"], section["content"],
                               theme_colors, is_sub=False)

        # 子节
        for sub in section["subsections"]:
            combined = section["content"] + sub["content"] if not section["content"] else sub["content"]
            if combined:
                title = f"{section['title']} — {sub['title']}" if sub["title"] else section["title"]
                _add_content_slide(prs, title, combined, theme_colors, is_sub=True)

    # 保存
    prs.save(output_path)
    file_size = os.path.getsize(output_path)
    return f"✅ PPT 已生成: {output_path} ({file_size / 1024:.0f}KB, {len(doc['sections'])} 章节)"


def main():
    parser = argparse.ArgumentParser(description="Markdown → PPTX 生成器")
    parser.add_argument("--input", "-i", required=True, help="输入 markdown 文件路径")
    parser.add_argument("--output", "-o", required=True, help="输出 .pptx 文件路径")
    parser.add_argument("--theme", default=DEFAULT_THEME,
                        choices=list(THEMES.keys()), help="配色方案")
    parser.add_argument("--brand-name", default="", help="品牌名（封面副标题）")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"❌ 输入文件不存在: {args.input}", file=sys.stderr)
        sys.exit(1)

    result = generate_ppt(args.input, args.output, args.theme, args.brand_name)
    print(result)


if __name__ == "__main__":
    main()
