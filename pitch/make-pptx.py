# -*- coding: utf-8 -*-
"""页图 → .pptx 组装（16:9 全幅贴图，页标题进演讲者备注）。

前置：node pitch/make-pptx.cjs 已产出 /tmp/pw/pptx-slides/slide-NN.png + titles.json
用法：python3 pitch/make-pptx.py
产出：pitch/AI养老院院长-产品介绍.pptx（图片页——版式/字体/图表与网页逐像素一致，
文字不可编辑；要改内容改 deck.html 后重跑两个脚本）
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SLIDES = Path("/tmp/pw/pptx-slides")
OUT = Path(__file__).parent / "AI养老院院长-产品介绍.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

titles = json.loads((SLIDES / "titles.json").read_text(encoding="utf-8"))
pngs = sorted(SLIDES.glob("slide-*.png"))
assert len(pngs) == len(titles), f"页数不一致: {len(pngs)} 图 vs {len(titles)} 标题"

for png, title in zip(pngs, titles):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    if title:
        slide.notes_slide.notes_text_frame.text = title

prs.save(OUT)
print(f"{OUT}: {OUT.stat().st_size / 1024 / 1024:.2f} MB, {len(pngs)} slides")
